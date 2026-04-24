"""
CourseForge PPT Builder
基于 python-pptx 生成 PPT，支持自定义模板
"""

import json
import os
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ────────────────────────────────────────
# 模板路径常量
# ────────────────────────────────────────
TEMPLATE_DIR = os.path.join(os.path.expanduser("~"), ".courseforge")
TEMPLATE_PATH = os.path.join(TEMPLATE_DIR, "template.pptx")


def get_saved_template_path() -> Optional[str]:
    """返回已保存的模板路径（如果存在）"""
    if os.path.isfile(TEMPLATE_PATH):
        return TEMPLATE_PATH
    return None


def save_template(file_bytes: bytes) -> str:
    """
    将上传的模板保存到 ~/.courseforge/template.pptx

    Args:
        file_bytes: 模板文件的字节内容

    Returns:
        保存后的文件路径
    """
    os.makedirs(TEMPLATE_DIR, exist_ok=True)
    with open(TEMPLATE_PATH, "wb") as f:
        f.write(file_bytes)
    return TEMPLATE_PATH


def delete_template() -> bool:
    """删除已保存的模板，返回是否成功"""
    if os.path.isfile(TEMPLATE_PATH):
        os.remove(TEMPLATE_PATH)
        return True
    return False


class PPTBuilder:

    @staticmethod
    def _build_chapter_index(prs, ai_generator) -> Dict[int, Dict[str, str]]:
        """
        [Phase 1] 全局章节扫描：提取幻灯片文本交由大模型，生成全局 Slide -> 章节映射
        """
        slide_contents = []
        # 扫描所有页面以实现全局聚类，通过字符截断防止 Token 超限
        for idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            joined_text = " ".join(texts)[:200]
            if joined_text:
                slide_contents.append(f"Slide {idx + 1}:\n{joined_text}")

        full_context = "\n\n".join(slide_contents)
        
        prompt = f"""
请分析以下 PPT 幻灯片的页面文本内容，识别出这套课件的完整「章节(Chapter)」和「小节(Section)」架构。
然后，为存在内容的幻灯片映射它所属的章节和小节。

文本内容：
{full_context}

【核心约束：强制粗粒度聚类（最少 5 页一节）】
为了保证后续题库生成的知识点密度，你必须对幻灯片进行强制的合并归类！
- 绝对底线：一个 `section`（小节）**至少必须包含 5 页连续的幻灯片**。绝对禁止划分出只有 1-4 页的细碎小节。
- 强制合并：请无视幻灯片之间细微的主题转换，强行将相邻的 5 页（或以上）幻灯片归入完全相同的 `chapter` 和 `section` 中。
- 示例要求：Slide 1 到 Slide 5 的 `chapter` 和 `section` 字段值必须一字不差地完全相同；Slide 6 到 Slide 10 同理。

【名称净素化规则】
- 绝对禁止在 `chapter` 和 `section` 的值中保留“第一章”、“第一部分”、“1.1”、“01”等一切数字或结构序号前缀！
- 必须只输出纯粹的名称文本。例如：将“第一章 规范”清洗为“规范”，将“1.1 礼仪”清洗为“礼仪”。

请严格输出以下 JSON 格式：
```json
{{
  "1": {{"chapter": "规范", "section": "礼仪"}},
  "2": {{"chapter": "规范", "section": "礼仪"}}
}}
```
注意：如果某页只是封面或目录，也尽量归入一个总括性的章/节。
必须仅输出 JSON，不要任何分析。
"""
        index_map = {}
        try:
            print("[VIBE-DEBUG] Phase 1: 正在生成全局章节映射...")
            response = ai_generator._call_ai(prompt, strip_thinking=True)
            import re, json
            match = re.search(r'```json\s*(.*?)\s*```', response, re.DOTALL)
            parsed_json = json.loads(match.group(1)) if match else json.loads(response)
            
            # 将字符串的 "1", "2" 转为 int 索引 (0-based: slide 1 -> index 0)
            for k, v in parsed_json.items():
                if str(k).isdigit():
                    index_map[int(k) - 1] = {
                        "chapter": v.get("chapter", "未知章节"),
                        "section": v.get("section", "未知小节")
                    }
        except Exception as e:
            print(f"[VIBE-DEBUG] Phase 1 章节映射生成失败，降级为逐页推断: {e}")
        
        return index_map

    @staticmethod
    def enrich_pptx_in_place(input_pptx_path: str, output_pptx_path: str, ai_generator, progress_callback=None) -> None:
        """
        从原有内容就地强化幻灯片备注：两步走 (Two-Pass) 策略
        1. 先抽部分文字生成全局树 (Chapter Index)
        2. 再逐页生成 Knowledge Base 并挂载回备注
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(input_pptx_path)
        
        # [Phase 1] 建立全局章节索引
        if progress_callback:
            progress_callback(0.05, "预读全文大纲，规划智能聚类 (Phase 1)...")
        chapter_index = PPTBuilder._build_chapter_index(prs, ai_generator)
        
        # [Phase 2] 逐页生成内容并写入
        print("[VIBE-DEBUG] Phase 2: 开始逐页萃取知识点...")
        total_slides = len(prs.slides)
        
        # 记录上一页的章节和小节，用于填补空白页
        last_chapter = "未分配章节"
        last_section = "未分配小节"

        for idx, slide in enumerate(prs.slides):
            if progress_callback:
                progress_callback(0.1 + 0.4 * (idx / max(total_slides, 1)), f"正在精读讲解第 {idx+1}/{total_slides} 页幻灯片 (Phase 2)...")
                
            slide_texts = []
            image_blob = None
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not image_blob:
                    try:
                        image_blob = shape.image.blob
                    except Exception:
                        pass
            
            text_content = "\n".join(slide_texts)
            
            # 无论文字多少，只要不是完全空白的骨架页，就进行处理
            # 此处拿掉原来的 > 50 字符限制，以便支持全图型 PPT
            if text_content.strip() or image_blob:
                try:
                    # 调用多模态底座
                    extracted = ai_generator.extract_slide_knowledge(text_content, image_blob=image_blob)
                    
                    # 合并全局索引里的 chapter / section（优先级：全局索引 > 逐页AI结果 > Last Known）
                    current_idx_meta = chapter_index.get(idx, {})
                    final_chapter = current_idx_meta.get("chapter") or extracted.get("chapter") or last_chapter
                    final_section = current_idx_meta.get("section") or extracted.get("section") or last_section
                    
                    # 只有当AI提取到了有意义的新章节名，且不是降级返回的"未知"，才更新last_known
                    # 为了防止 "待补充" 或 "未知" 污染上下文
                    if final_chapter and final_chapter not in ["未知", "待补充", "未知章节"]:
                        last_chapter = final_chapter
                    if final_section and final_section not in ["未知", "待补充", "未知小节"]:
                        last_section = final_section
                    
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    
                    final_text = f"[{last_chapter}][{last_section}]\n{extracted.get('knowledge_base', '')}"
                    text_frame.text = final_text
                    print(f"  -> Slide {idx+1} 处理完成: [{last_chapter}][{last_section}]")
                except Exception as e:
                    print(f"[VIBE-DEBUG] Slide {idx+1} 知识点提炼失败: {e}")
                    
        prs.save(output_pptx_path)

    """PPT 构建器 — 支持自定义模板"""

    # 默认配色方案（极简商务风，无模板时使用）
    COLOR_SCHEME = {
        'primary': RGBColor(44, 62, 80),
        'secondary': RGBColor(52, 73, 94),
        'text': RGBColor(44, 62, 80),
        'accent': RGBColor(41, 128, 185),
        'background': RGBColor(255, 255, 255)
    }

    # 版式名称关键字 → 内部类型映射
    _LAYOUT_KEYWORDS = {
        'title': ['标题幻灯片', 'Title Slide', '封面'],
        'section': ['节标题', 'Section Header', '章节', 'Section'],
        'content': ['标题和内容', 'Title and Content', '内容', 'Content'],
        'blank': ['空白', 'Blank'],
    }

    def __init__(self, template_path: Optional[str] = None):
        """
        初始化 PPT 构建器
        
        Args:
            template_path: 可选，.pptx 模板文件路径。
        """
        self._use_template = False
        self._layout_map: Dict[str, Any] = {}
        self.prs = None

        # 尝试加载模板 (Safe Template Mode)
        if template_path and os.path.isfile(template_path):
            try:
                # 1. 加载模板
                temp_prs = Presentation(template_path)
                
                # 2. 预检测版式 (借用 _detect_layouts 逻辑，但针对 temp_prs)
                temp_map = self._scan_layouts(temp_prs)
                
                # 3. 严格校验 (Validate)
                # 必须包含: 封面, 章节, 内容页
                required_keys = {'title', 'section', 'content'}
                if required_keys.issubset(temp_map.keys()):
                    #进一步检查内容页是否有 2 个以上占位符 (Title + Content)
                    content_layout = temp_map['content']
                    if len(content_layout.placeholders) >= 2:
                        self.prs = temp_prs
                        self._layout_map = temp_map
                        self._use_template = True
                        self._clear_existing_slides()
                        print(f"[PPTBuilder] Template loaded successfully: {template_path}")
            except Exception as e:
                print(f"[PPTBuilder] Template validation failed: {str(e)}")
                self.prs = None

        # Fallback: 如果模板无效或未提供，使用默认空白母版
        if not self.prs:
            self.prs = Presentation()
            self._use_template = False
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
            print("[PPTBuilder] Using default code-drawn style (Safe Mode)")

    def _clear_existing_slides(self) -> None:
        """清除模板中已有的幻灯片（保留母版和版式）"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].get('r:id')
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def _scan_layouts(self, prs_obj) -> Dict[str, Any]:
        """
        扫描指定 Presentation 对象的版式
        """
        layout_map = {}
        for layout in prs_obj.slide_layouts:
            name = layout.name
            matched = False
            for type_key, keywords in self._LAYOUT_KEYWORDS.items():
                for kw in keywords:
                    # 排除竖向版式
                    if 'vertical' in name.lower() or '竖' in name:
                        continue
                        
                    if kw.lower() in name.lower():
                        if type_key not in layout_map:
                            layout_map[type_key] = layout
                        matched = True
                        break
                if matched:
                    break
        return layout_map

    def _detect_layouts(self) -> None:
        """
        (Legacy) 兼容旧调用，扫描当前 self.prs
        """
        self._layout_map = self._scan_layouts(self.prs)

    def _get_layout(self, slide_type: str):
        """
        获取指定类型的版式，优先使用模板版式，否则 fallback 到空白。

        Args:
            slide_type: 'title' | 'section' | 'content' | 'blank'
        """
        if self._use_template and slide_type in self._layout_map:
            return self._layout_map[slide_type]
        
        # Fallback: 必须使用真正的空白版式（无占位符），以触发 _fallback_xxx 及其绘图逻辑
        # 默认 python-pptx 模板中，index 6 是空白页
        blank_layout = None
        for i, layout in enumerate(self.prs.slide_layouts):
            if "blank" in layout.name.lower() or "空白" in layout.name:
                blank_layout = layout
                break
        
        if blank_layout:
            return blank_layout
            
        # 如果找不到名字叫 blank 的，尝试使用 index 6 (通常是空白)
        if len(self.prs.slide_layouts) > 6:
            return self.prs.slide_layouts[6]
            
        # 实在没办法，用最后一个
        return self.prs.slide_layouts[len(self.prs.slide_layouts) - 1]

    def _try_set_placeholder(self, slide, idx: int, text: str,
                              font_size: int = None, bold: bool = None,
                              color: RGBColor = None,
                              alignment=None) -> bool:
        """
        尝试向 placeholder 写入文本。成功返回 True，不存在返回 False。
        当使用模板时，只写入文本内容，不覆盖任何格式属性（完全继承母版样式）。

        Args:
            slide: 幻灯片对象
            idx: placeholder 索引
            text: 要写入的文本
            font_size: 可选字号（模板模式下忽略）
            bold: 可选加粗（模板模式下忽略）
            color: 可选颜色（模板模式下忽略）
            alignment: 可选对齐（模板模式下忽略）
        """
        try:
            ph = slide.placeholders[idx]
            ph.text = text
            # 无论是否使用模板，都强制应用极简风格样式（用户要求）
            if font_size:
                ph.text_frame.paragraphs[0].font.size = Pt(font_size)
            if bold is not None:
                ph.text_frame.paragraphs[0].font.bold = bold
            if color:
                ph.text_frame.paragraphs[0].font.color.rgb = color
            if alignment:
                ph.text_frame.paragraphs[0].alignment = alignment
            return True
        except (KeyError, IndexError):
            return False

    def _add_title_slide(self, title: str, subtitle: str = "") -> None:
        """添加标题页"""
        layout = self._get_layout('title')
        slide = self.prs.slides.add_slide(layout)

        # 尝试使用模板 placeholder
        title_set = self._try_set_placeholder(
            slide, 0, title,
            font_size=44, bold=True,
            color=self.COLOR_SCHEME['primary'],
            alignment=PP_ALIGN.CENTER
        )
        if subtitle:
            self._try_set_placeholder(
                slide, 1, subtitle,
                font_size=24,
                color=self.COLOR_SCHEME['secondary'],
                alignment=PP_ALIGN.CENTER
            )

        # Fallback: 模板没有 placeholder 时手动添加
        if not title_set:
            self._fallback_title_slide(slide, title, subtitle)

    def _fallback_title_slide(self, slide, title: str, subtitle: str) -> None:
        """标题页 fallback（手动 textbox）"""
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.COLOR_SCHEME['primary']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(0.8)
            )
            sf = sub_box.text_frame
            sf.text = subtitle
            sf.paragraphs[0].font.size = Pt(24)
            sf.paragraphs[0].font.color.rgb = self.COLOR_SCHEME['secondary']
            sf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, title: str) -> None:
        """添加章节分隔页"""
        layout = self._get_layout('section')
        slide = self.prs.slides.add_slide(layout)

        title_set = self._try_set_placeholder(
            slide, 0, title,
            font_size=40, bold=True,
            color=self.COLOR_SCHEME['primary']
        )

        if not title_set:
            self._fallback_section_slide(slide, title)

    def _fallback_section_slide(self, slide, title: str) -> None:
        """章节页 fallback"""
        slide.shapes.add_shape(
            1, Inches(0.5), Inches(2),
            Inches(0.15), Inches(3.5),
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = self.COLOR_SCHEME['accent']

        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3), Inches(7), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.COLOR_SCHEME['primary']

    def _add_content_slide(self, title: str, bullets: List[str]) -> None:
        """添加内容页"""
        layout = self._get_layout('content')
        slide = self.prs.slides.add_slide(layout)

        # 标题
        title_set = self._try_set_placeholder(
            slide, 0, title,
            font_size=32, bold=True,
            color=self.COLOR_SCHEME['primary']
        )

        # 内容（placeholder[1] = body）
        body_set = False
        try:
            body_ph = slide.placeholders[1]
            body_ph.text = ""
            tf = body_ph.text_frame
            tf.clear()
            for i, bullet_text in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
                # 强制应用样式
                p.font.size = Pt(20)
                p.font.color.rgb = self.COLOR_SCHEME['text']
                p.space_before = Pt(12) if i > 0 else Pt(0)
            body_set = True
        except (KeyError, IndexError):
            pass

        # Fallback
        if not title_set or not body_set:
            if not title_set:
                self._fallback_content_title(slide, title)
            if not body_set:
                self._fallback_content_body(slide, bullets)

    def _fallback_content_title(self, slide, title: str) -> None:
        """内容页标题 fallback"""
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.COLOR_SCHEME['primary']

        slide.shapes.add_connector(
            1, Inches(0.8), Inches(1.5), Inches(9.2), Inches(1.5)
        ).line.color.rgb = self.COLOR_SCHEME['accent']

    def _fallback_content_body(self, slide, bullets: List[str]) -> None:
        """内容页正文 fallback"""
        content_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2), Inches(7.6), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            p.text = f"• {bullet_text}"
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLOR_SCHEME['text']
            p.space_before = Pt(12) if i > 0 else Pt(0)

    def _add_end_slide(self, title: str = "感谢聆听", subtitle: str = "Q&A") -> None:
        """添加结束页"""
        layout = self._get_layout('title')  # 结束页复用标题版式
        slide = self.prs.slides.add_slide(layout)

        title_set = self._try_set_placeholder(
            slide, 0, title,
            font_size=48, bold=True,
            color=self.COLOR_SCHEME['primary'],
            alignment=PP_ALIGN.CENTER
        )
        if subtitle:
            self._try_set_placeholder(
                slide, 1, subtitle,
                font_size=28,
                color=self.COLOR_SCHEME['accent'],
                alignment=PP_ALIGN.CENTER
            )

        if not title_set:
            self._fallback_title_slide(slide, title, subtitle)

    def build_from_json(self, json_outline: str) -> Presentation:
        """
        从 JSON 大纲构建 PPT

        Args:
            json_outline: JSON 格式的 PPT 大纲

        Returns:
            生成的 Presentation 对象
        """
        try:
            json_text = json_outline.strip()

            # 提取 JSON
            start_idx = json_text.find('{')
            end_idx = json_text.rfind('}')

            if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
                json_text = json_text[start_idx: end_idx + 1]
            else:
                if "```json" in json_text:
                    parts = json_text.split("```json")
                    if len(parts) > 1:
                        json_text = parts[1].split("```")[0]
                elif "```" in json_text:
                    parts = json_text.split("```")
                    if len(parts) > 1:
                        json_text = parts[1]

            json_text = json_text.strip()
            outline = json.loads(json_text)

            for slide_def in outline.get('slides', []):
                slide_type = slide_def.get('type', 'content')

                if slide_type == 'title':
                    self._add_title_slide(
                        slide_def.get('title', '课程标题'),
                        slide_def.get('subtitle', '')
                    )
                elif slide_type == 'section':
                    self._add_section_slide(slide_def.get('title', '章节标题'))
                elif slide_type == 'content':
                    self._add_content_slide(
                        slide_def.get('title', '内容标题'),
                        slide_def.get('bullets', [])
                    )
                elif slide_type == 'end':
                    self._add_end_slide(
                        slide_def.get('title', '感谢聆听'),
                        slide_def.get('subtitle', 'Q&A')
                    )

            return self.prs

        except json.JSONDecodeError as e:
            snippet = json_text[:200] + "..." if len(json_text) > 200 else json_text
            raise ValueError(
                f"JSON 解析失败: {str(e)}\n\n提取后的文本片段:\n{snippet}\n\n原始文本全部长度: {len(json_outline)}"
            )
        except Exception as e:
            raise RuntimeError(f"PPT 生成失败: {str(e)}")

    def save(self, filepath: str) -> None:
        """保存 PPT 文件"""
        self.prs.save(filepath)


# ────────────────────────────────────────
# 便捷函数
# ────────────────────────────────────────
def generate_ppt_from_outline(
    json_outline: str,
    output_path: str,
    template_path: Optional[str] = None
) -> str:
    """
    从 JSON 大纲生成 PPT 文件

    Args:
        json_outline: JSON 格式的大纲
        output_path: 输出文件路径
        template_path: 可选，.pptx 模板路径

    Returns:
        生成的文件路径
    """
    builder = PPTBuilder(template_path=template_path)
    builder.build_from_json(json_outline)
    builder.save(output_path)
    return output_path
