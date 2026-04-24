"""
CourseForge File Parser
支持 PDF, DOCX, PPTX 格式的课件解析
"""

import io
from typing import Optional
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


class FileParser:
    """统一的文件解析器"""
    
    @staticmethod
    def parse_pdf(file_bytes: bytes) -> str:
        """
        解析 PDF 文件
        
        Args:
            file_bytes: PDF 文件的字节流
            
        Returns:
            提取的文本内容
        """
        try:
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            
            text_content = []
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text_content.append(f"--- 第 {page_num} 页 ---\n{page_text}\n")
            
            return "\n".join(text_content)
        
        except Exception as e:
            raise ValueError(f"PDF 解析失败: {str(e)}")
    
    @staticmethod
    def parse_docx(file_bytes: bytes) -> str:
        """
        解析 DOCX 文件
        
        Args:
            file_bytes: DOCX 文件的字节流
            
        Returns:
            提取的文本内容
        """
        try:
            docx_file = io.BytesIO(file_bytes)
            doc = Document(docx_file)
            
            text_content = []
            
            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    # 识别标题
                    if para.style.name.startswith('Heading'):
                        text_content.append(f"\n## {para.text}\n")
                    else:
                        text_content.append(para.text)
            
            # 提取表格
            for table_num, table in enumerate(doc.tables, 1):
                text_content.append(f"\n--- 表格 {table_num} ---")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    text_content.append(row_text)
                text_content.append("")
            
            return "\n".join(text_content)
        
        except Exception as e:
            raise ValueError(f"DOCX 解析失败: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_bytes: bytes) -> str:
        """
        解析 PPTX 文件
        
        Args:
            file_bytes: PPTX 文件的字节流
            
        Returns:
            提取的文本内容
        """
        try:
            pptx_file = io.BytesIO(file_bytes)
            prs = Presentation(pptx_file)
            
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content.append(f"\n{'='*50}")
                text_content.append(f"幻灯片 {slide_num}")
                text_content.append('='*50)
                
                # 提取所有文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 尝试识别标题
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            text_content.append(f"\n[标题] {shape.text}")
                        else:
                            text_content.append(shape.text)
                    
                    # 提取表格
                    if shape.has_table:
                        text_content.append("\n[表格]")
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            text_content.append(row_text)
                
                text_content.append("")
            
            return "\n".join(text_content)
        
        except Exception as e:
            raise ValueError(f"PPTX 解析失败: {str(e)}")
    
    @classmethod
    def parse_file(cls, file_bytes: bytes, file_type: str) -> str:
        """
        根据文件类型自动选择解析器
        
        Args:
            file_bytes: 文件字节流
            file_type: 文件类型 ('pdf', 'docx', 'pptx')
            
        Returns:
            解析后的文本内容
        """
        parsers = {
            'pdf': cls.parse_pdf,
            'docx': cls.parse_docx,
            'pptx': cls.parse_pptx
        }
        
        parser = parsers.get(file_type.lower())
        if not parser:
            raise ValueError(f"不支持的文件类型: {file_type}")
        
        content = parser(file_bytes)
        
        # 基础清洗：去除多余空行
        lines = [line for line in content.split('\n') if line.strip() or line == '']
        cleaned_content = '\n'.join(lines)
        
        return cleaned_content
    
    @staticmethod
    def clean_content(text: str) -> str:
        """
        深度清洗文本内容（去除 LOGO、公司名等）
        
        Args:
            text: 原始文本
            
        Returns:
            清洗后的文本
        """
        # 常见需要过滤的关键词模式
        sensitive_keywords = [
            'copyright', '©', '版权所有',
            'logo', 'LOGO',
            '公司', '集团', '有限责任公司', 'Co., Ltd',
            '保密', '内部资料'
        ]
        
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # 跳过包含敏感词的行
            if any(keyword.lower() in line.lower() for keyword in sensitive_keywords):
                continue
            cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)


# 便捷函数
def extract_text_from_file(uploaded_file) -> tuple[str, str]:
    """
    从 Streamlit 上传的文件中提取文本
    
    Args:
        uploaded_file: Streamlit UploadedFile 对象
        
    Returns:
        (原始内容, 清洗后的内容)
    """
    file_bytes = uploaded_file.read()
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    # 解析文件
    raw_content = FileParser.parse_file(file_bytes, file_extension)
    
    # 清洗内容
    cleaned_content = FileParser.clean_content(raw_content)
    
    return raw_content, cleaned_content
