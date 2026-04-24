"""
CourseForge (铸课工坊) - 本地化 AI 备课助手
主应用程序 v2.0 - 支持程序内配置
"""

import streamlit as st
import os
from datetime import datetime
from modules import (
    extract_text_from_file,
    create_generator,
    generate_ppt_from_outline,
    ConfigManager,
    load_or_create_config,
    is_configured,
    get_system_prompt,
    save_template,
    delete_template,
    get_saved_template_path
)


# ==================== 页面配置 ====================
st.set_page_config(
    page_title="CourseForge 铸课工坊",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)


# ==================== 样式定义 ====================
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #2c3e50;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        font-size: 1.2rem;
        color: #7f8c8d;
        text-align: center;
        margin-bottom: 2rem;
    }
    .step-box {
        background-color: #ecf0f1;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #3498db;
        margin-bottom: 1rem;
    }
    .success-box {
        background-color: #d4edda;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #28a745;
        margin: 1rem 0;
    }
    .warning-box {
        background-color: #fff3cd;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #ffc107;
        margin: 1rem 0;
    }
    .config-box {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 0.5rem;
        border: 1px solid #dee2e6;
        margin: 1rem 0;
    }

    /* Hide Streamlit components - TEMPORARILY DISABLED FOR DEBUGGING
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    [data-testid="stToolbar"] {visibility: hidden;}
    .stDeployButton {display:none;}
    */
</style>
""", unsafe_allow_html=True)


# ==================== 初始化 Session State ====================
def init_session_state():
    """初始化会话状态"""
    if 'page' not in st.session_state:
        st.session_state.page = 'check_config'
    if 'config' not in st.session_state:
        st.session_state.config = load_or_create_config()
    if 'show_config' not in st.session_state:
        st.session_state.show_config = False


# ==================== 配置页面 ====================
def show_config_page():
    """显示配置页面"""
    st.markdown('<div class="main-header">⚙️ API 配置</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">首次使用需要配置 API 信息</div>', unsafe_allow_html=True)
    
    # ========== API 模式选择（在 form 外部） ==========
    current_mode = st.session_state.config.get('api_mode', 'external')
    mode_index = 0 if current_mode == 'internal' else 1
    
    api_mode_label = st.radio(
        "🌐 选择 API 模式",
        options=["🏢 内网模式（Internal AI Platform）", "🌍 外网模式（Google / Anthropic）"],
        index=mode_index,
        horizontal=True,
        help="内网模式使用企业内部 AI 平台的 appkey 鉴权；外网模式使用 Google GenAI / Anthropic SDK"
    )
    api_mode = 'internal' if '内网' in api_mode_label else 'external'
    
    # ========== 根据模式显示不同的说明 ==========
    if api_mode == 'internal':
        st.markdown("""
        <div class="config-box">
        <h4>📖 内网模式配置说明</h4>
        <p>连接企业内部 AI 平台，数据不出域，安全合规。</p>
        <p><b>如何获取配置？</b></p>
        <ul>
            <li>登录「内部控制台」获取接口地址</li>
            <li>在「资源中心」→「详情」中复制「应用标识」(appkey)</li>
        </ul>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <div class="config-box">
        <h4>📖 外网模式配置说明</h4>
        <p>CourseForge 需要连接到兼容 OpenAI API 的服务（如 Gemini）来生成教学材料。</p>
        <p><b>如何获取 API？</b></p>
        <ul>
            <li>推荐：Google AI Studio (Gemini API) - 免费额度充足</li>
            <li>备选：企业内部 AI 服务、其他兼容平台</li>
        </ul>
        </div>
        """, unsafe_allow_html=True)
    
    # ========== 配置表单 ==========
    with st.form("config_form"):
        
        if api_mode == 'internal':
            # ---------- 内网模式 ----------
            st.subheader("🔑 内部 AI 平台配置")
            
            appkey = st.text_input(
                "应用标识 (appkey) *",
                value=st.session_state.config.get('appkey', ''),
                type="password",
                placeholder="YOUR_APP_KEY_HERE",
                help="从「内部控制台」资源中心获取的应用标识"
            )
            
            internal_url = st.text_input(
                "接口地址 (POST URL) *",
                value=st.session_state.config.get('internal_url', ''),
                placeholder="https://your-internal-api.com/v1",
                help="内部大模型接口的完整 POST 地址"
            )
            
            # 此处填写贵司内部大模型通道
            internal_models = [
                "internal-model-fast",    # 快推理模型
                "internal-model-deep",    # 深度推理模型
                "internal-model-max",     # 满血版模型
                "internal-model-long",    # 长文本模型
            ]
            
            current_internal_model = st.session_state.config.get('internal_model', 'internal-model-fast')
            internal_model_index = internal_models.index(current_internal_model) if current_internal_model in internal_models else 0
            
            internal_model = st.selectbox(
                "模型通道 *",
                options=internal_models,
                index=internal_model_index,
                help="此处选择内部配置的模型通道"
            )
            
            # 将内网参数映射到通用变量名，方便后续统一处理
            api_key = ""
            base_url = internal_url
            model_name = internal_model
        else:
            # ---------- 外网模式 ----------
            st.subheader("🔑 外网 API 配置")
            
            api_key = st.text_input(
                "API Key *",
                value=st.session_state.config.get('api_key', ''),
                type="password",
                placeholder="sk-xxxxxxxxxxxxxxxx",
                help="从 API 服务商获取的密钥"
            )
            
            base_url = st.text_input(
                "Base URL (选填)",
                value=st.session_state.config.get('base_url', 'http://127.0.0.1:8045'),
                placeholder="直连 Google API 可留空或填 https://generativelanguage.googleapis.com",
                help="外部代理或中转 API 地址（留空默认直连 Google）"
            )

            is_google_direct = (not base_url.strip()) or ("googleapis.com" in base_url)
            
            if is_google_direct:
                model_options = [
                    "gemini-3.1-pro-preview",
                    "gemini-3.1-flash-lite-preview",
                    "gemini-2.5-flash"
                ]
                help_text = "选择 Google 官方模型"
            else:
                model_options = [
                    "gemini-3.1-pro-high",
                    "claude-opus-4-6-thinking"
                ]
                help_text = "选择 Antigravity 支持的模型"

            # 尝试恢复之前保存的配置模型，如果不在列表里则选第一个
            saved_model = st.session_state.config.get('model_name', model_options[0])
            try:
                default_index = model_options.index(saved_model)
            except ValueError:
                default_index = 0
                
            model_name = st.selectbox(
                "模型选择 *",
                options=model_options,
                index=default_index,
                help=help_text
            )
            
            appkey = ""
            internal_url = ""
            internal_model = ""
        
        st.divider()
        
        st.subheader("🎨 自定义系统 Prompt（可选）")
        
        use_custom_prompt = st.checkbox(
            "使用自定义系统 Prompt",
            value=bool(st.session_state.config.get('custom_system_prompt')),
            help="如不勾选，将使用默认的动态行业模板"
        )
        
        custom_prompt = ""
        if use_custom_prompt:
            default_prompt = get_system_prompt("通用企业")

            custom_prompt = st.text_area(
                "系统 Prompt",
                value=st.session_state.config.get('custom_system_prompt', default_prompt),
                height=300,
                help="定义 AI 的角色、约束和任务流程"
            )
        
        col1, col2 = st.columns(2)
        
        with col1:
            submit = st.form_submit_button("💾 保存配置", type="primary", use_container_width=True)
        
        with col2:
            test = st.form_submit_button("🧪 测试连接", use_container_width=True)
        
        if submit:
            # 构建配置（包含两种模式的所有字段，确保切换时不丢失数据）
            new_config = {
                'api_mode': api_mode,
                'api_key': api_key or st.session_state.config.get('api_key', ''),
                'base_url': base_url if api_mode == 'external' else st.session_state.config.get('base_url', ''),
                'model_name': model_name if api_mode == 'external' else st.session_state.config.get('model_name', 'gemini-2.5-pro'),
                'custom_system_prompt': custom_prompt if use_custom_prompt else '',
                'appkey': appkey or st.session_state.config.get('appkey', ''),
                'internal_url': internal_url or st.session_state.config.get('internal_url', ''),
                'internal_model': internal_model or st.session_state.config.get('internal_model', 'internal-model-fast')
            }
            
            # 验证配置
            is_valid, error_msg = ConfigManager.validate_config(new_config)
            
            if not is_valid:
                st.error(f"❌ 配置验证失败: {error_msg}")
            else:
                # 保存配置
                if ConfigManager.save_config(new_config):
                    st.session_state.config = new_config
                    st.success("✅ 配置保存成功！")
                    st.balloons()
                    st.session_state.page = 'main'
                    st.rerun()
                else:
                    st.error("❌ 配置保存失败，请检查文件权限")
        
        if test:
            # 测试 API 连接
            if api_mode == 'internal':
                if not appkey or not internal_url:
                    st.warning("⚠️ 请先填写应用标识和接口地址")
                else:
                    with st.spinner("🔌 正在测试内部大模型连接..."):
                        try:
                            test_generator = create_generator(
                                api_key="",
                                base_url=internal_url,
                                model_name=internal_model,
                                custom_system_prompt=custom_prompt if use_custom_prompt else None,
                                api_mode='internal',
                                appkey=appkey
                            )
                            
                            test_response = test_generator._call_ai(
                                "请回复：连接成功",
                                temperature=0.1,
                                max_tokens=50
                            )
                            
                            st.success(f"✅ 内部大模型连接成功！\n\n模型响应: {test_response}")
                        
                        except Exception as e:
                            st.error(f"❌ 连接失败: {str(e)}\n\n请检查：\n1. appkey（应用标识）是否正确\n2. 接口地址是否正确\n3. 内网网络是否连通")
            else:
                if not api_key:
                    st.warning("⚠️ 请先填写 API Key")
                else:
                    with st.spinner("🔌 正在测试 API 连接..."):
                        try:
                            test_generator = create_generator(
                                api_key, 
                                base_url, 
                                model_name,
                                custom_prompt if use_custom_prompt else None,
                                api_mode='external'
                            )
                            
                            # 发送测试请求
                            test_response = test_generator._call_ai(
                                "请回复：连接成功",
                                temperature=0.1,
                                max_tokens=50
                            )
                            
                            st.success(f"✅ API 连接成功！\n\n模型响应: {test_response}")
                        
                        except Exception as e:
                            st.error(f"❌ API 连接失败: {str(e)}\n\n请检查：\n1. API Key 是否正确\n2. Base URL 是否正确\n3. 网络连接是否正常")


# ==================== 主界面 ====================
def show_main_page():
    """显示主界面"""
    
    # 标题
    st.markdown('<div class="main-header">📚 CourseForge 铸课工坊</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">AI 驱动的一键备课助手 | 本地部署 · 隐私安全</div>', unsafe_allow_html=True)

    # 添加右上角配置按钮作为备选入口
    col_header_1, col_header_2 = st.columns([6, 1])
    with col_header_2:
        if st.button("⚙️ API 设置", key="main_config_btn", help="点击修改 API 配置"):
            st.session_state.show_config = True
            st.session_state.page = 'config'
            st.rerun()

    # 获取当前配置
    config = st.session_state.config
    
    # 侧边栏 - 配置区
    with st.sidebar:
        st.header("⚙️ 配置")

        # 显示当前配置详情
        api_mode = config.get('api_mode', 'external')
        
        if api_mode == 'internal':
            # 内网模式显示
            appkey_val = config.get('appkey', '')
            if len(appkey_val) > 8:
                appkey_masked = f"{appkey_val[:4]}...{appkey_val[-4:]}"
            else:
                appkey_masked = "未配置"
            
            st.markdown(f"""
            <div class="step-box">
            <b>当前配置状态</b> <span style="color: #27ae60;">🏢 内网模式</span><br>
            <small>🤖 模型: {config.get('internal_model', 'internal-model-fast')}</small><br>
            <small>🔗 地址: {config.get('internal_url', '未配置')}</small><br>
            <small>🔑 appkey: {appkey_masked}</small>
            </div>
            """, unsafe_allow_html=True)
        else:
            # 外网模式显示
            api_key = config.get('api_key', '')
            if len(api_key) > 8:
                api_key_masked = f"{api_key[:3]}...{api_key[-4:]}"
            else:
                api_key_masked = "未配置"

            st.markdown(f"""
            <div class="step-box">
            <b>当前配置状态</b> <span style="color: #2980b9;">🌍 外网模式</span><br>
            <small>🤖 模型: {config.get('model_name', 'gemini-2.5-pro')}</small><br>
            <small>🔗 地址: {config.get('base_url', '未配置')}</small><br>
            <small>🔑 Key: {api_key_masked}</small>
            </div>
            """, unsafe_allow_html=True)
        
        # 配置管理
        if st.button("🔧 修改配置", use_container_width=True):
            st.session_state.show_config = True
            st.session_state.page = 'config'
            st.rerun()
        
        if st.button("🗑️ 删除配置", use_container_width=True):
            if ConfigManager.delete_config():
                st.success("配置已删除")
                st.session_state.config = ConfigManager.get_default_config()
                st.session_state.page = 'config'
                st.rerun()
        
        st.divider()
        
        # 行业场景选择
        industry = st.radio(
            "🏢 行业场景",
            options=["通用企业", "银行业", "互联网"],
            help="不同行业会使用对应的术语和业务场景"
        )

        st.divider()
        
        # 受众级别选择
        audience_level = st.radio(
            "🎯 受众级别",
            options=["初级/基石营", "中级/先锋营"],
            help="不同级别会调整内容的深度和风格"
        )
        
        st.divider()
        
        # 生成选项
        st.subheader("📋 生成内容")
        
        st.caption("🎓 先导系列")
        generate_precourse = st.checkbox("先导课程大纲", value=True)
        generate_quiz = st.checkbox("场景化通关测试", value=True)
        
        st.caption("📚 核心系列")
        generate_video = st.checkbox("视频旁白脚本", value=True)
        generate_interaction = st.checkbox("课堂互动方案", value=True)
        generate_action = st.checkbox("回岗实践计划", value=True)
        generate_survey = st.checkbox("双向调研问卷", value=True)
        generate_ppt = st.checkbox("标准化 PPT (重新生成)", value=True)

        st.divider()

        st.subheader("🔬 级联强化 (LMS模式)")
        enable_lms = st.checkbox("启用 LMS 闭环处理", value=True, help="限PPTX上传：1.保留原版式插入AI强化备注; 2.生成配套LMS题库Excel")

        st.divider()
        
        # PPT 模板设置
        st.subheader("🎨 PPT 模板")
        
        saved_tpl = get_saved_template_path()
        if saved_tpl:
            st.success("✅ 已配置企业模板")
            col_tpl1, col_tpl2 = st.columns(2)
            with col_tpl1:
                if st.button("🔄 更换模板", use_container_width=True, key="replace_tpl"):
                    st.session_state['show_tpl_uploader'] = True
            with col_tpl2:
                if st.button("🗑️ 删除模板", use_container_width=True, key="del_tpl"):
                    delete_template()
                    st.success("模板已删除，将使用默认风格")
                    st.rerun()
        else:
            st.info("📎 使用默认极简风格")
            st.session_state.setdefault('show_tpl_uploader', True)
        
        if st.session_state.get('show_tpl_uploader', not saved_tpl):
            tpl_file = st.file_uploader(
                "上传 PPT 模板 (.pptx)",
                type=['pptx'],
                help="上传企业标准 PPT 模板，生成的课件将自动套用模板的母版、配色和字体",
                key="tpl_uploader"
            )
            if tpl_file:
                save_template(tpl_file.getvalue())
                st.success(f"✅ 模板已保存: {tpl_file.name}")
                st.session_state['show_tpl_uploader'] = False
                st.rerun()
        
        st.divider()
        
        # 使用说明
        with st.expander("📖 使用说明"):
            st.markdown("""
            **操作步骤：**
            1. 上传原始课件（PDF/DOCX/PPTX）
            2. 选择受众级别
            3. 勾选需要生成的内容
            4. 点击"一键生成"
            
            **注意事项：**
            - 文件大小建议 < 10MB
            - 生成时间约 2-5 分钟
            - 所有数据均在本地处理
            - AI 生成内容仅供参考
            """)
    
    # 主内容区
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.subheader("📤 上传课件")
        uploaded_file = st.file_uploader(
            "支持 PDF, DOCX, PPTX 格式",
            type=['pdf', 'docx', 'pptx'],
            help="请上传需要转化的原始课件"
        )
        
        if uploaded_file:
            st.success(f"✅ 已上传: {uploaded_file.name} ({uploaded_file.size / 1024:.1f} KB)")
            
            # 预览原始内容
            with st.expander("🔍 预览原始内容"):
                try:
                    raw_content, cleaned_content = extract_text_from_file(uploaded_file)
                    st.text_area(
                        "提取的文本（前 1000 字）",
                        cleaned_content[:1000] + "...",
                        height=300
                    )
                except Exception as e:
                    st.error(f"文件解析失败: {str(e)}")
    
    with col2:
        st.subheader("🎨 课程信息")
        course_name = st.text_input(
            "课程名称",
            placeholder="例如：高效沟通技巧",
            help="用于文件命名"
        )
    
    # 生成按钮
    st.divider()
    
    if st.button("🚀 一键生成全套教学材料", type="primary", use_container_width=True):
        
        # 验证输入
        if not uploaded_file:
            st.error("❌ 请先上传课件文件")
            return
        
        if not course_name.strip():
            # 使用上传文件名（去扩展名）作为默认课程名称
            course_name = os.path.splitext(uploaded_file.name)[0]
        
        # 开始生成
        try:
            # 提取文件内容
            with st.spinner("📖 正在解析课件..."):
                # 重置文件指针
                uploaded_file.seek(0)
                raw_content, cleaned_content = extract_text_from_file(uploaded_file)
            
            st.success("✅ 文件解析完成")
            
            # 创建 AI 生成器（根据 api_mode 自动路由）
            api_mode = config.get('api_mode', 'external')
            if api_mode == 'internal':
                generator = create_generator(
                    api_key='',
                    base_url=config.get('internal_url', ''),
                    model_name=config.get('internal_model', 'internal-model-fast'),
                    custom_system_prompt=config.get('custom_system_prompt') or None,
                    api_mode='internal',
                    appkey=config.get('appkey', ''),
                    industry=industry
                )
            else:
                generator = create_generator(
                    config['api_key'],
                    config['base_url'],
                    config['model_name'],
                    config.get('custom_system_prompt') or None,
                    api_mode='external',
                    industry=industry
                )
            
            # 进度条
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            def update_progress(message: str, progress: float):
                status_text.text(message)
                progress_bar.progress(progress)
            
            # 收集生成选项
            options = {
                'precourse_outline': generate_precourse,
                'scenario_quiz': generate_quiz,
                'video_script': generate_video,
                'interactions': generate_interaction,
                'action_plan': generate_action,
                'surveys': generate_survey,
                'ppt_outline': generate_ppt,
                'enable_lms': enable_lms
            }

            # 生成内容
            with st.spinner("🤖 AI 正在生成教学材料..."):
                results = generator.generate_all_materials(
                    cleaned_content,
                    audience_level,
                    options=options,
                    progress_callback=update_progress,
                    course_name=course_name
                )
            
            # 保存文件
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_dir = f"outputs/{course_name}_{timestamp}"
            os.makedirs(output_dir, exist_ok=True)
            
            saved_files = {}
            
            # 保存文本文件
            text_outputs = {
                'core_content.md': results.get('core_content', ''),
                'precourse_outline.md': results.get('precourse_outline', ''),
                'video_script.md': results.get('video_script', ''),
                'interactions.md': results.get('interactions', ''),
                'action_plan.md': results.get('action_plan', ''),
                'surveys.md': results.get('surveys', ''),
                'scenario_quiz.md': results.get('scenario_quiz', {}).get('markdown', '') if isinstance(results.get('scenario_quiz'), dict) else results.get('scenario_quiz', '')
            }
            
            for filename, content in text_outputs.items():
                if content:
                    filepath = os.path.join(output_dir, filename)
                    with open(filepath, 'w', encoding='utf-8') as f:
                        f.write(content)
                    saved_files[filename] = filepath
            
            # 生成并保存 PPT（自动套用模板）
            if results.get('ppt_outline'):
                try:
                    ppt_path = os.path.join(output_dir, f"{course_name}_课件.pptx")
                    tpl_path = get_saved_template_path()
                    generate_ppt_from_outline(results['ppt_outline'], ppt_path, template_path=tpl_path)
                    saved_files['课件.pptx'] = ppt_path
                except Exception as e:
                    st.warning(f"PPT 生成失败: {str(e)}")
                    
            # ========= 注入 LMS 处理逻辑 =============
            lms_pptx_path = None
            lms_excel_path = None
            if options.get('enable_lms') and uploaded_file.name.lower().endswith(".pptx"):
                try:
                    from modules.ppt_builder import PPTBuilder
                    from modules.excel_exporter import ExcelExporter
                    from modules.prompts import PromptTemplates
                    import json, re
                    
                    update_progress("🔥 正在进行 LMS 闭环处理（PPT备注强化 + 题库生成）...", 0.9)
                    st.session_state.ai_generator = generator
                    
                    # 1. 强化 PPT
                    temp_pptx = os.path.join(output_dir, "temp_source.pptx")
                    lms_pptx_path = os.path.join(output_dir, f"{course_name}_LMS强化备注.pptx")
                    
                    uploaded_file.seek(0)
                    with open(temp_pptx, "wb") as f:
                        f.write(uploaded_file.read())
                        
                    PPTBuilder.enrich_pptx_in_place(temp_pptx, lms_pptx_path, generator)
                    saved_files['LMS强化备注.pptx'] = lms_pptx_path
                    
                    # 2. 导出 Excel
                    lms_excel_path = os.path.join(output_dir, f"{course_name}_LMS题库.xlsx")
                    
                    # 提取强化后的 PPT 备注，以便大模型能读取到 [章节][小节] 标签
                    try:
                        from pptx import Presentation
                        prs_enriched = Presentation(lms_pptx_path)
                        notes_texts = []
                        for slide in prs_enriched.slides:
                            # 提取备注
                            note_text = ""
                            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                                note_text = slide.notes_slide.notes_text_frame.text.strip()
                            # 提取正文兜底
                            shape_texts = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    shape_texts.append(shape.text.strip())
                            
                            combined_text = note_text
                            if shape_texts and note_text:
                                combined_text += "\n【正文原文】\n" + "\n".join(shape_texts)
                            elif shape_texts:
                                combined_text = "\n".join(shape_texts)
                                
                            if combined_text:
                                notes_texts.append(combined_text)
                        
                        ppt_full_text = "\n\n---\n\n".join(notes_texts)
                        if not ppt_full_text.strip():
                            ppt_full_text = cleaned_content
                    except Exception as e:
                        print(f"读取强化版备注失败: {e}")
                        ppt_full_text = cleaned_content
                        
                    # 将 ppt_full_text 按 [章节][小节] 标签进行分块 (Chunking) 以防大模型偷懒或触发 Max Tokens 截断
                    import re
                    chunks = []
                    current_chunk = []
                    current_label = "默认章节"
                    
                    for line in ppt_full_text.split('\n'):
                        # 匹配类似 [第一章][1.1 概念] 或单独的 [第一章]
                        match = re.match(r'^\[(.*?)\](?:\[(.*?)\])?', line.strip())
                        if match:
                            if current_chunk:
                                chunks.append((current_label, "\n".join(current_chunk)))
                                current_chunk = []
                            current_label = line.strip()
                            current_chunk.append(line)
                        else:
                            current_chunk.append(line)
                    if current_chunk:
                        chunks.append((current_label, "\n".join(current_chunk)))
                    
                    all_questions = []
                    
                    # 针对没有明显分块的短 PPT，兜底当做一块处理
                    if len(chunks) == 0:
                        chunks = [("全局内容", ppt_full_text)]
                        
                    progress_step = 0.9 / max(len(chunks), 1)
                    current_progress = 0.0
                    
                    for i, (label, chunk_text) in enumerate(chunks):
                        if not chunk_text.strip(): continue
                        msg = f"⏳ 正在为 {label} 生成强制配额题库 ({i+1}/{len(chunks)})..."
                        update_progress(msg, current_progress)
                        
                        prompt = PromptTemplates.get_scenario_quiz_prompt(chunk_text, audience_level, course_name)
                        try:
                            ai_response = generator._call_ai(prompt, max_tokens=4000, strip_thinking=True)
                            chunk_questions = generator.extract_json_from_response(ai_response)
                            if chunk_questions:
                                all_questions.extend(chunk_questions)
                        except Exception as e:
                            print(f"[VIBE-DEBUG] Chunk {label} 生成题库失败: {e}")
                            
                        current_progress += progress_step
                        
                    questions = all_questions
                             
                    # 防御性检查：如果解析不出题目，跳过 Excel 导出
                    if not questions:
                        st.warning("⚠️ AI 返回了内容，但未能解析出结构化题目 JSON。请检查模型输出格式。")
                        st.text_area("AI 原始响应（调试用）", ai_response[:2000], height=200)
                    else:
                        update_progress(f"📊 正在导出 {len(questions)} 道题目到 Excel...", 0.95)
                        # 定位模板
                        import sys
                        if getattr(sys, 'frozen', False):
                            base_path = sys._MEIPASS
                        else:
                            base_path = os.path.dirname(os.path.abspath(__file__))
                            
                        template_excel = os.path.join(base_path, 'assets', '导入题库模板.xlsx')
                        if not os.path.exists(template_excel):
                            st.error("Excel 模板文件丢失，请确保 'assets/导入题库模板.xlsx' 存在！")
                            raise FileNotFoundError("模板文件丢失")
                        
                        try:
                            ExcelExporter.export_lms_questions(template_excel, lms_excel_path, questions, course_name=course_name)
                            saved_files['LMS题库.xlsx'] = lms_excel_path
                            st.success(f"✅ LMS 题库导出完成！共 {len(questions)} 道题目")
                        except PermissionError:
                            st.error("请先关闭 Excel 模板文件！")
                            raise
                        except Exception as e:
                            st.error(f"导出 Excel 失败: {str(e)}")
                            raise
                    
                except Exception as e:
                    import traceback
                    st.warning(f"⚠️ LMS 混合强化发生异常: {e}")
                    print(traceback.format_exc())
            
            # 清除进度条
            progress_bar.empty()
            status_text.empty()
            
            # 显示成功消息
            st.markdown(f"""
            <div class="success-box">
            <h3>🎉 生成完成！</h3>
            <p>所有文件已保存到: <code>{output_dir}</code></p>
            </div>
            """, unsafe_allow_html=True)
            
            # 显示生成结果
            st.subheader("📄 生成结果")
            
            tabs = st.tabs([
                "📝 核心内容",
                "🎓 先导大纲",
                "📝 通关测试",
                "🎬 视频脚本",
                "🎮 课堂互动",
                "🎯 回岗实践",
                "📊 调研问卷",
                "📑 PPT 大纲"
            ])
            
            with tabs[0]:
                if results.get('core_content'):
                    st.markdown(results['core_content'])
                    st.download_button(
                        "💾 下载 Markdown",
                        results['core_content'],
                        file_name="core_content.md",
                        mime="text/markdown"
                    )
            
            with tabs[1]:
                if generate_precourse and results.get('precourse_outline'):
                    st.markdown(results['precourse_outline'])
                    st.download_button(
                        "💾 下载先导大纲",
                        results['precourse_outline'],
                        file_name="precourse_outline.md",
                        mime="text/markdown"
                    )
            
            with tabs[2]:
                if generate_quiz and results.get('scenario_quiz'):
                    quiz_data = results['scenario_quiz']
                    quiz_markdown = quiz_data.get('markdown', str(quiz_data)) if isinstance(quiz_data, dict) else quiz_data
                    
                    st.markdown(quiz_markdown)
                    st.download_button(
                        "💾 下载通关测试",
                        quiz_markdown,
                        file_name="scenario_quiz.md",
                        mime="text/markdown"
                    )
            
            with tabs[3]:
                if generate_video and results.get('video_script'):
                    st.markdown(results['video_script'])
                    st.download_button(
                        "💾 下载 Markdown",
                        results['video_script'],
                        file_name="video_script.md",
                        mime="text/markdown"
                    )
            
            with tabs[4]:
                if generate_interaction and results.get('interactions'):
                    st.markdown(results['interactions'])
                    st.download_button(
                        "💾 下载 Markdown",
                        results['interactions'],
                        file_name="interactions.md",
                        mime="text/markdown"
                    )
            
            with tabs[5]:
                if generate_action and results.get('action_plan'):
                    st.markdown(results['action_plan'])
                    st.download_button(
                        "💾 下载 Markdown",
                        results['action_plan'],
                        file_name="action_plan.md",
                        mime="text/markdown"
                    )
            
            with tabs[6]:
                if generate_survey and results.get('surveys'):
                    st.markdown(results['surveys'])
                    st.download_button(
                        "💾 下载 Markdown",
                        results['surveys'],
                        file_name="surveys.md",
                        mime="text/markdown"
                    )
            
            with tabs[7]:
                if generate_ppt and results.get('ppt_outline'):
                    st.code(results['ppt_outline'], language='json')
                    
                    if saved_files.get('课件.pptx'):
                        with open(saved_files['课件.pptx'], 'rb') as f:
                            st.download_button(
                                "💾 下载 重新排版PPT",
                                f.read(),
                                file_name=f"{course_name}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                            
                    col_btn1, col_btn2 = st.columns(2)
                    with col_btn1:
                        if 'lms_pptx_path' in locals() and lms_pptx_path and os.path.exists(lms_pptx_path):
                            with open(lms_pptx_path, 'rb') as f:
                                st.download_button(
                                    "🚀 下载 .pptx (PPT 强化版)",
                                    f.read(),
                                    file_name=f"{course_name}_LMS强化备注.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                                
                    with col_btn2:
                        if 'lms_excel_path' in locals() and lms_excel_path and os.path.exists(lms_excel_path):
                            with open(lms_excel_path, 'rb') as f:
                                st.download_button(
                                    "📊 下载 .xlsx (Excel 题库版)",
                                    f.read(),
                                    file_name=f"{course_name}_LMS题库.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                )
        
        except Exception as e:
            st.error(f"❌ 生成失败: {str(e)}")
            st.exception(e)
    
    # ==================== 独立工具：MD 题库 → Excel 转换 ====================
    st.divider()
    st.subheader("🔄 独立工具：文档/场景化测试 → Excel 题库")
    st.caption("上传已生成的 `scenario_quiz.md` 或 自建的 `docx` 题库文档，独立转换为 Excel 题库。")
    
    with st.expander("👉 查看支持的题库排版格式 (为保证转写成功率，请尽量贴近此格式)"):
        st.markdown("""
        **支持的极简格式规范：**
        1. **题目序号**：使用数字加点开头，如 `1. ` 或 `2.`
        2. **章节标签**（可选）：紧跟序号后，使用方括号，如 `[第一章][1.1]`
        3. **题干**：标签后的所有文字
        4. **选项**：使用大写字母加点开头，如 `A. xxx` (仅单选/多选需要)
        5. **答案**：另起一行，以 `答案：` 或 `正确答案：` 开头
        6. **解析**（可选）：以 `解析：` 或 `深度解析：` 开头
        
        **示例**：
        ```text
        1. [第一章][1.1 概念] 以下哪项属于AI的特点？
        A. 智能计算
        B. 传统算盘
        C. 机械齿轮
        答案：A
        解析：AI具备智能计算能力。
        ```
        """)
    
    quiz_md_file = st.file_uploader(
        "上传 scenario_quiz.md, lms_content.md 或 .docx 题库",
        type=["md", "txt", "docx"],
        key="quiz_md_uploader",
        help="支持场景化通关测试 Markdown 格式 或 DOCX 格式"
    )
    
    if quiz_md_file is not None:
        file_ext = os.path.splitext(quiz_md_file.name)[1].lower()
        
        if file_ext == '.docx':
            from modules.file_parser import FileParser
            quiz_md_file.seek(0)
            md_content = FileParser.parse_docx(quiz_md_file.read())
        else:
            md_content = quiz_md_file.read().decode('utf-8')
        
        from modules.quiz_converter import QuizConverter
        
        # 先尝试混合解析（JSON 优先 + Markdown 回退）
        parsed_questions = QuizConverter.parse_lms_mixed_response(md_content)
        
        if not parsed_questions:
            st.error("❌ 未能从文件中解析出任何题目。请检查文件格式是否正确。")
            st.text_area("文件内容预览", md_content[:2000], height=200)
        else:
            st.success(f"✅ 成功解析 **{len(parsed_questions)}** 道题目！")
            
            # 预览前 3 题（含元数据）
            with st.expander("📋 题目预览（前 3 题）", expanded=True):
                for i, q in enumerate(parsed_questions[:3]):
                    meta_parts = []
                    if q.get('course_name'):
                        meta_parts.append(f"📚 {q['course_name']}")
                    if q.get('chapter'):
                        meta_parts.append(f"📖 {q['chapter']}")
                    if q.get('section'):
                        meta_parts.append(f"📄 {q['section']}")
                    meta_line = " | ".join(meta_parts) if meta_parts else ""
                    
                    if meta_line:
                        st.caption(meta_line)
                    st.markdown(f"**{i+1}. {q.get('question', '(无题干)')[:120]}...**")
                    for opt in q.get('options', []):
                        prefix = "✅" if opt.get('is_correct') else "⬜"
                        st.markdown(f"  {prefix} {opt.get('text', '')}")
                    st.markdown("---")
            
            # 自动提取课程名称
            auto_course_name = parsed_questions[0].get('course_name', '') if parsed_questions else ''
            convert_course_name = st.text_input("课程名称（用于 Excel A 列，已自动提取）", value=auto_course_name or "微课模块", key="convert_course_name")
            
            if st.button("📊 立即转换为 Excel 题库", key="btn_convert_excel"):
                import sys
                if getattr(sys, 'frozen', False):
                    base_path = sys._MEIPASS
                else:
                    base_path = os.path.dirname(os.path.abspath(__file__))
                    
                template_excel = os.path.join(base_path, 'assets', '导入题库模板.xlsx')
                
                if not os.path.exists(template_excel):
                    st.error("Excel 模板文件丢失，请确保 'assets/导入题库模板.xlsx' 存在！")
                else:
                    try:
                        from modules.excel_exporter import ExcelExporter
                        
                        output_excel = os.path.join(os.path.expanduser("~"), "Desktop", f"{convert_course_name}_LMS题库.xlsx")
                        ExcelExporter.export_lms_questions(template_excel, output_excel, parsed_questions, course_name=convert_course_name)
                        
                        st.success(f"🎉 导出成功！文件已保存到：`{output_excel}`")
                        
                        with open(output_excel, 'rb') as f:
                            st.download_button(
                                "⬇️ 下载 Excel 题库",
                                f.read(),
                                file_name=f"{convert_course_name}_LMS题库.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            )
                    except Exception as e:
                        st.error(f"❌ 导出失败: {str(e)}")
                        import traceback
                        print(traceback.format_exc())
    
    # 页脚
    st.divider()
    st.markdown("""
    <div style="text-align: center; color: #7f8c8d; font-size: 0.9rem;">
        <p>CourseForge v2.8 | 本地部署 · 数据安全 · 开源免费</p>
        <p>⚠️ AI 生成内容仅供参考，请人工审核后使用</p>
    </div>
    """, unsafe_allow_html=True)


# ==================== 主程序逻辑 ====================
def main():
    """主应用逻辑"""
    init_session_state()
    
    # 检查配置状态
    if st.session_state.page == 'check_config':
        if is_configured():
            st.session_state.page = 'main'
        else:
            st.session_state.page = 'config'
    
    # 路由到对应页面
    if st.session_state.page == 'config':
        show_config_page()
    else:
        show_main_page()


if __name__ == "__main__":
    main()



# ----------------------------------------------------
# 核心业务流：打通 LMS 闭环 (PPT强化 + 考题生成 + Excel装填)
# ----------------------------------------------------
def process_course_pipeline(uploaded_pptx, template_excel_path):
    st.info("🚀 正在并行执行核心业务流：就地挂载能力、生成先导大纲及试题...")
    
    output_enriched_pptx = "temp_Enriched_Course.pptx"
    output_excel_path = "temp_LMS_Questions.xlsx"
    temp_pptx_path = "temp_source.pptx"
    
    with open(temp_pptx_path, "wb") as f:
        f.write(uploaded_pptx.read())

    # [Task 1] & [Task 2]
    progress_bar = st.progress(0.0)
    status_text = st.empty()
    
    def update_progress(ratio, msg):
        progress_bar.progress(ratio)
        status_text.info(f"⏳ {msg}")

    with st.spinner("流水线运算中，这可能需要几分钟，请勿刷新页面..."):
        # 1. 触发备注强化 (内部自行处理图片和多模态)
        from modules.ppt_builder import PPTBuilder
        PPTBuilder.enrich_pptx_in_place(temp_pptx_path, output_enriched_pptx, st.session_state.ai_generator, progress_callback=update_progress)
        
        update_progress(0.55, "幻灯片智能重构完毕。正在提炼结构化记忆节点...")
        # 2. 生成题目和微课大纲
        # 此处不再使用 extract_text_for_ai 提取原始内容
        # 而是直接读取刚才强化好的 PPT 备注（里面包含了 [章节][小节] 和原文+解读）
        try:
            from pptx import Presentation
            prs_enriched = Presentation(output_enriched_pptx)
            notes_texts = []
            for slide in prs_enriched.slides:
                if slide.has_notes_slide:
                    text_frame = slide.notes_slide.notes_text_frame
                    if text_frame.text.strip():
                        notes_texts.append(text_frame.text.strip())
            ppt_full_text = "\n\n---\n\n".join(notes_texts)
            if not ppt_full_text.strip():
                raise ValueError("Enriched notes are empty")
        except Exception as e:
            print(f"[VIBE-DEBUG] 读取强化版备注失败: {e}，降级")
            ppt_full_text = "请基于通用情境生成。"
            
        update_progress(0.60, f"提取到大纲基础文本（长达 {len(ppt_full_text)} 字）。正在指派 AI 考评专家出题（由于配额限制，此过程耗时约为 2~5 分钟...）")
        from modules.prompts import PromptTemplates
        prompt = PromptTemplates.get_lms_content_prompt(ppt_full_text)
        # 扩大 token 上限，防止 10 题 * N 小节导致 JSON 截断
        ai_response = st.session_state.ai_generator._call_ai(prompt, max_tokens=8192, strip_thinking=True)
        update_progress(0.85, "AI 题库生成完毕，正在进行降噪校验与格式装配...")

    # 拆分 Markdown 和 JSON
    import json, re
    import re
    json_str = "[]"
    markdown_content = ai_response
    json_match = re.search(r'```json\s*(.*?)\s*```', ai_response, re.DOTALL)
    if json_match:
        markdown_content = ai_response[:json_match.start()].strip()
        json_str = json_match.group(1).strip()

    # UI 预览大纲
    st.markdown("### 微课大纲与题目预览")
    st.markdown(markdown_content)

    # [Task 4] - JSON与Excel防御解析
    questions = []
    try:
        from modules.ai_generator import AIGenerator
        questions = AIGenerator.extract_json_from_response(ai_response)
        update_progress(0.95, f"成功解析 {len(questions)} 道题目，准备导出专属试卷格式...")
    except Exception as e:
        import logging
        logging.error(f"⚠️ [VIBE-DEBUG] AI 生成题库结果解析异常: {e}")
        st.error(f"⚠️ [VIBE-DEBUG] AI 生成题库结果解析异常: {e}")

    try:
        # [Task 3] 填充Excel
        from modules.excel_exporter import ExcelExporter
        template_excel = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets', '导入题库模板.xlsx')
        if not os.path.exists(template_excel):
            st.error("Excel 模板文件丢失，请确保 'assets/导入题库模板.xlsx' 存在！")
        else:
            ExcelExporter.export_lms_questions(template_excel, output_excel_path, questions)
            
            # UI：多轨并行下载按钮
            st.success("✅ 闭环处理完成，请下载所需资源：")
        col1, col2 = st.columns(2)
        
        with col1:
            with open(output_enriched_pptx, "rb") as f:
                st.download_button(
                    label="⬇️ 下载 .pptx (PPT 强化版)",
                    data=f.read(),
                    file_name="PPT_Enriched_Notes.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                
        with col2:
            with open(output_excel_path, "rb") as f:
                st.download_button(
                    label="⬇️ 下载 .xlsx (Excel 题库版)",
                    data=f.read(),
                    file_name="LMS_Questions_Export.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
    except PermissionError:
        st.error("❌ 请先关闭 Excel 模板文件！")
    except Exception as e:
        st.error(f"❌ 发生了未知导出错误：{e}")
